from fastapi import FastAPI, Request, HTTPException, UploadFile, File
from fastapi.responses import HTMLResponse, FileResponse
from anthropic import Anthropic
import os
import re
import json
import base64                   
from io import BytesIO
from dotenv import load_dotenv
import docx  # python-docx
from pptx import Presentation  # python-pptx
from pypdf import PdfReader  # pypdf
from PIL import Image
from xhtml2pdf import pisa
from bs4 import BeautifulSoup

load_dotenv()

client = Anthropic(
    api_key=os.getenv("ANTHROPIC_API_KEY"),
    base_url=os.getenv("ANTHROPIC_BASE_URL"),
    default_headers={
        "anthropic-workspace-id": os.getenv("ANTHROPIC_WORKSPACE_ID")
    }
)

app = FastAPI()

SYSTEM_PROMPT = """
You are a helpful assistant capable of generating text, documents, slides, and graphics.
If the user requests an asset, format it inside code blocks exactly like this:
1. IMAGES/GRAPHICS: Pure valid SVG code inside ```xml
2. PDF Documents: Complete, beautifully styled HTML page inside ```html
3. WORD/DOCX Documents: Organized Markdown structure inside ```markdown
4. POWERPOINT/PPTX Slides: Structured JSON payload inside ```json representing individual slides. Example:
[
  {"title": "Slide Title", "content": ["Bullet point 1", "Bullet point 2"]},
  {"title": "Second Slide", "content": ["Point A", "Point B"]}
]

Always wrap structural code outputs inside these code blocks so the parser cleanly intercepts them.
"""

history = []
pending_attachments = []

# --- MARKDOWN TO HTML CONVERTER FOR PDF EXPORTS ---
def process_inline_markdown(text):
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    text = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"__(.*?)__", r"<strong>\1</strong>", text)
    text = re.sub(r"`(.*?)`", r"<code>\1</code>", text)
    return text

def clean_markdown_symbols(text):
    return text.replace("**", "").replace("`", "").replace("__", "")

def markdown_to_html(md_text):
    html = """<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
    @page {
        size: letter;
        margin: 1in;
    }
    body {
        font-family: Helvetica, Arial, sans-serif;
        color: #2d3748;
        line-height: 1.6;
        font-size: 11pt;
    }
    #footer {
        position: fixed;
        bottom: -0.5in;
        right: 0;
        font-size: 8pt;
        color: #888888;
        text-align: right;
    }
    h1 {
        font-size: 22pt;
        color: #1a365d;
        border-bottom: 2px solid #e2e8f0;
        padding-bottom: 6px;
        margin-top: 0;
        margin-bottom: 12px;
    }
    h2 {
        font-size: 16pt;
        color: #2b6cb0;
        margin-top: 18px;
        margin-bottom: 8px;
    }
    h3 {
        font-size: 12pt;
        color: #4a5568;
        margin-top: 14px;
        margin-bottom: 6px;
    }
    p {
        margin-bottom: 10px;
    }
    ul, ol {
        margin-bottom: 12px;
        padding-left: 20px;
    }
    li {
        margin-bottom: 4px;
    }
    pre {
        background-color: #f7fafc;
        border: 1px solid #e2e8f0;
        border-radius: 4px;
        padding: 8px;
        font-family: Courier, monospace;
        font-size: 9.5pt;
        margin-bottom: 12px;
    }
    code {
        font-family: Courier, monospace;
        background-color: #f7fafc;
        padding: 2px 4px;
        border-radius: 3px;
        font-size: 9.5pt;
    }
    strong, b {
        font-weight: bold;
        color: #1a202c;
    }
</style>
</head>
<body>
<div id="footer">Page <pdf:pagenumber> of <pdf:pagecount></div>
"""
    lines = md_text.split("\n")
    in_code = False
    code_content = []
    in_list = False
    
    for line in lines:
        stripped = line.strip()
        
        if stripped.startswith("```"):
            if in_code:
                in_code = False
                html += "<pre><code>" + "".join(code_content) + "</code></pre>\n"
                code_content = []
            else:
                in_code = True
            continue
            
        if in_code:
            escaped = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            code_content.append(escaped + "\n")
            continue
            
        if stripped.startswith("* ") or stripped.startswith("- "):
            if not in_list:
                in_list = True
                html += "<ul>\n"
            content = stripped[2:]
            html += f"<li>{process_inline_markdown(content)}</li>\n"
            continue
        else:
            if in_list:
                in_list = False
                html += "</ul>\n"
                
        if stripped.startswith("# "):
            html += f"<h1>{process_inline_markdown(stripped[2:])}</h1>\n"
        elif stripped.startswith("## "):
            html += f"<h2>{process_inline_markdown(stripped[3:])}</h2>\n"
        elif stripped.startswith("### "):
            html += f"<h3>{process_inline_markdown(stripped[4:])}</h3>\n"
        elif stripped:
            html += f"<p>{process_inline_markdown(stripped)}</p>\n"
            
    if in_list:
        html += "</ul>\n"
        
    html += "</body></html>"
    return html


# --- PDF GENERATOR HELPERS ---
def clean_and_fix_html(html_str):
    # If style block is open, close it
    if re.search(r"<style", html_str, re.IGNORECASE) and not re.search(r"</style>", html_str, re.IGNORECASE):
        html_str += "\n</style>"
        
    # Find style blocks and fix their CSS braces
    style_matches = re.findall(r"(<style[\s\S]*?>)([\s\S]*?)(</style>)", html_str, re.IGNORECASE)
    for match in style_matches:
        full_style = match[0] + match[1] + match[2]
        css_text = match[1]
        
        # fix braces in css_text
        open_c = css_text.count("{")
        close_c = css_text.count("}")
        if open_c > close_c:
            fixed_css = css_text.rstrip()
            if not fixed_css.endswith(";") and not fixed_css.endswith("}"):
                fixed_css += ";"
            fixed_css += "\n}" * (open_c - close_c)
            fixed_style = match[0] + fixed_css + match[2]
            html_str = html_str.replace(full_style, fixed_style)
            
    # Parse with BeautifulSoup to auto-close other HTML tags
    soup = BeautifulSoup(html_str, "html.parser")
    return str(soup)

def convert_html_to_pdf(html_content, output_pdf_path):
    fixed_html = clean_and_fix_html(html_content)
    with open(output_pdf_path, "wb") as pdf_file:
        pisa_status = pisa.CreatePDF(fixed_html, dest=pdf_file)
    return not pisa_status.err


# --- DOCX GENERATOR HELPERS ---
def add_markdown_paragraph(doc, line):
    is_list = False
    if line.startswith("* ") or line.startswith("- "):
        is_list = True
        line = line[2:]
    elif line.startswith("1. "):
        is_list = True
        line = line[3:]
        
    style = 'List Bullet' if is_list else None
    p = doc.add_paragraph(style=style) if style else doc.add_paragraph()
    
    parts = re.split(r"(\*\*.*?\*\*)", line)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            p.add_run(part[2:-2]).bold = True
        else:
            p.add_run(part)



# --- DOCX PARSER & GENERATOR HELPERS ---
def write_markdown_to_docx(doc, markdown_content):
    lines = markdown_content.split("\n")
    i = 0
    num_lines = len(lines)
    
    while i < num_lines:
        line = lines[i]
        stripped = line.strip()
        
        if not stripped:
            i += 1
            continue
            
        # Check if table
        if stripped.startswith("|") and i + 1 < num_lines and re.match(r"^\|[\s\-\|:]+\|$", lines[i+1].strip()):
            table_rows = []
            # Header cells
            table_rows.append([cell.strip() for cell in stripped.split("|")[1:-1]])
            i += 2 # Skip separator
            
            while i < num_lines and lines[i].strip().startswith("|"):
                row_cells = [cell.strip() for cell in lines[i].strip().split("|")[1:-1]]
                table_rows.append(row_cells)
                i += 1
                
            if table_rows:
                num_rows = len(table_rows)
                num_cols = max(len(row) for row in table_rows)
                table = doc.add_table(rows=num_rows, cols=num_cols)
                table.style = 'Table Grid'
                
                for r_idx, row_data in enumerate(table_rows):
                    row = table.rows[r_idx]
                    for c_idx, cell_value in enumerate(row_data):
                        if c_idx < len(row.cells):
                            cell = row.cells[c_idx]
                            p = cell.paragraphs[0]
                            run = p.add_run(cell_value)
                            if r_idx == 0:
                                run.bold = True
            continue
            
        if stripped.startswith("# "):
            doc.add_heading(clean_markdown_symbols(stripped[2:]), level=1)
            i += 1
        elif stripped.startswith("## "):
            doc.add_heading(clean_markdown_symbols(stripped[3:]), level=2)
            i += 1
        elif stripped.startswith("### "):
            doc.add_heading(clean_markdown_symbols(stripped[4:]), level=3)
            i += 1
        else:
            add_markdown_paragraph(doc, stripped)
            i += 1



# --- INGESTION HELPERS FOR DOCX & PPTX ---
def parse_docx_file(file_bytes):
    doc = docx.Document(BytesIO(file_bytes))
    content = []
    
    # Extract headers
    for section in doc.sections:
        if section.header and section.header.text.strip():
            content.append(f"[Header: {section.header.text.strip()}]")
            
    # Iterate through body elements in order of appearance
    for child in doc.element.body:
        if child.tag.endswith('p'):
            p = docx.text.paragraph.Paragraph(child, doc)
            if p.text.strip():
                content.append(p.text.strip())
        elif child.tag.endswith('tbl'):
            table = docx.table.Table(child, doc)
            table_lines = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                # Filter out adjacent duplicate cells that python-docx gets from merged cells
                cleaned_cells = []
                for cell in cells:
                    if not cleaned_cells or cell != cleaned_cells[-1]:
                        cleaned_cells.append(cell)
                table_lines.append(" | ".join(cleaned_cells))
            if table_lines:
                content.append("\n[Table]\n" + "\n".join(table_lines) + "\n[/Table]")
                
    # Extract footers
    for section in doc.sections:
        if section.footer and section.footer.text.strip():
            content.append(f"[Footer: {section.footer.text.strip()}]")
            
    return "\n".join(content)

def parse_pptx_file(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        slide_texts.append(f"=== Slide {i+1} ===")
        # Extract title shape if present
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_texts.append(f"Title: {slide.shapes.title.text.strip()}")
            
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            # Handle text frame shapes
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            # Handle tables in slides
            if shape.has_table:
                table_lines = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_lines.append(" | ".join(cells))
                slide_texts.append("Table:\n" + "\n".join(table_lines))
                
        # Extract notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"Notes: {notes}")
                
    return "\n".join(slide_texts)


def parse_excel_file(file_bytes):
    import openpyxl
    wb = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True)
    sheets_content = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheets_content.append(f"=== Sheet: {sheet_name} ===")
        
        # Read rows
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
            
        # Determine non-empty rows
        non_empty_rows = []
        for row in rows:
            if any(cell is not None for cell in row):
                non_empty_rows.append(row)
                
        if not non_empty_rows:
            continue
            
        # Find maximum length of row to size columns properly
        max_cols = max(len(row) for row in non_empty_rows)
        
        # Convert rows to Markdown table format
        for idx, row in enumerate(non_empty_rows):
            # Pad row values if some rows have fewer elements
            row_cells = list(row) + [None] * (max_cols - len(row))
            cleaned_row = [str(cell).replace("\n", " ").strip() if cell is not None else "" for cell in row_cells]
            sheets_content.append("| " + " | ".join(cleaned_row) + " |")
            if idx == 0:
                separator = "|" + "|".join(["---" for _ in cleaned_row]) + "|"
                sheets_content.append(separator)
        sheets_content.append("") # empty line after sheet
        
    return "\n".join(sheets_content)


# --- PPTX PARSER HELPERS ---
def parse_markdown_to_slides(text):
    slides = []
    lines = text.split("\n")
    current_slide = None
    
    for line in lines:
        line_str = line.strip()
        if not line_str:
            continue
            
        if line_str.startswith("# ") or line_str.startswith("## ") or line_str.startswith("### "):
            if current_slide:
                slides.append(current_slide)
            title = line_str.lstrip("#").strip()
            current_slide = {"title": title, "content": []}
        elif current_slide is not None:
            if line_str.startswith("- ") or line_str.startswith("* ") or line_str.startswith("1. "):
                bullet_content = re.sub(r"^(\-\s*|\*\s*|\d+\.\s*)", "", line_str)
                current_slide["content"].append(bullet_content)
            else:
                current_slide["content"].append(line_str)
        else:
            current_slide = {"title": "Summary", "content": [line_str]}
            
    if current_slide:
        slides.append(current_slide)
        
    return [s for s in slides if s["title"] or s["content"]]

def parse_html_to_slides(html_str):
    soup = BeautifulSoup(html_str, "html.parser")
    slides = []
    
    # Filter class names to match exact target page containers only
    pages = soup.find_all(class_=lambda x: x and any(cls in ["cover", "page", "slide"] for cls in x.split()))
    
    for page in pages:
        title_el = page.find(["h1", "h2", "h3"])
        if not title_el:
            title_el = page.find(class_=re.compile(r"(title|header)", re.IGNORECASE))
        title = title_el.get_text().strip() if title_el else "Slide"
        
        content = []
        for li in page.find_all("li"):
            txt = li.get_text().strip()
            if txt:
                content.append(txt)
                
        card_classes = ["card", "stat-box", "box", "ai-type", "app-card", "stat-box-alt", "stat-box-green", "stat-box-orange", "quote-block"]
        for card in page.find_all(class_=lambda x: x and any(cls in card_classes for cls in x.split())):
            card_title_el = card.find(["h3", "h4", "h5"])
            card_title = card_title_el.get_text().strip() if card_title_el else ""
            
            card_desc_el = card.find("p")
            card_desc = card_desc_el.get_text().strip() if card_desc_el else ""
            
            if card_title and card_desc:
                content.append(f"{card_title}: {card_desc}")
            elif card_title:
                content.append(card_title)
            elif card_desc:
                content.append(card_desc)
            else:
                txt = card.get_text().strip()
                if txt:
                    content.append(" ".join(txt.split()))
        
        if not content:
            for p in page.find_all("p"):
                if any(cls in p.get("class", []) for cls in ["subtitle", "footer", "page-number"]):
                    continue
                txt = p.get_text().strip()
                if txt and len(txt) < 300:
                    content.append(txt)
                    
        slides.append({"title": title, "content": content})
        
    if not slides:
        current_slide = None
        for element in soup.find_all(["h1", "h2", "h3", "p", "li"]):
            if element.name in ["h1", "h2", "h3"]:
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": element.get_text().strip(), "content": []}
            elif current_slide:
                txt = element.get_text().strip()
                if txt:
                    current_slide["content"].append(txt)
        if current_slide:
            slides.append(current_slide)
            
    return slides

def extract_slides_from_text(text):
    # Check for JSON code block first
    json_match = re.search(r"```json\s*([\s\S]*?)(?:```|$)", text, re.IGNORECASE)
    if json_match:
        try:
            return json.loads(json_match.group(1).strip())
        except Exception:
            pass
            
    # Try parsing text directly as JSON
    try:
        return json.loads(text.strip())
    except Exception:
        pass

    html_match = re.search(r"```html\s*([\s\S]*?)(?:```|$)", text, re.IGNORECASE)
    if html_match:
        html_content = html_match.group(1).strip()
        return parse_html_to_slides(html_content)
    elif "<html>" in text.lower() or "<!doctype html>" in text.lower():
        return parse_html_to_slides(text)
    else:
        return parse_markdown_to_slides(text)


@app.get("/", response_class=HTMLResponse)
async def home():
    with open("index.html", "r", encoding="utf-8") as f:
        return f.read()


# --- FILE UPLOAD INGESTION ENDPOINT ---
@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    filename = file.filename.lower()
    file_bytes = await file.read()
    global pending_attachments
    
    try:
        # 1. Process Word files (.docx)
        if filename.endswith(".docx"):
            extracted_text = parse_docx_file(file_bytes)
            pending_attachments.append({
                "filename": file.filename,
                "type": "text",
                "content": extracted_text
            })
            return {
                "status": f"Parsed DOCX content from {file.filename}",
                "attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]
            }

        # 2. Process PowerPoint files (.pptx)
        elif filename.endswith(".pptx"):
            extracted_text = parse_pptx_file(file_bytes)
            pending_attachments.append({
                "filename": file.filename,
                "type": "text",
                "content": extracted_text
            })
            return {
                "status": f"Parsed PPTX content from {file.filename}",
                "attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]
            }

        # 3. Process PDF files (.pdf)
        elif filename.endswith(".pdf"):
            reader = PdfReader(BytesIO(file_bytes))
            pdf_texts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pdf_texts.append(text)
            extracted_text = "\n".join(pdf_texts)
            encoded_pdf = base64.b64encode(file_bytes).decode("utf-8")
            pending_attachments.append({
                "filename": file.filename,
                "type": "pdf",
                "media_type": "application/pdf",
                "data": encoded_pdf,
                "content": extracted_text
            })
            return {
                "status": f"Attached PDF content from {file.filename}",
                "attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]
            }

        # 3b. Process Excel files (.xlsx, .xls)
        elif filename.endswith((".xlsx", ".xls")):
            extracted_text = parse_excel_file(file_bytes)
            pending_attachments.append({
                "filename": file.filename,
                "type": "text",
                "content": extracted_text
            })
            return {
                "status": f"Parsed Excel content from {file.filename}",
                "attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]
            }

        # 4. Process Images (supporting png, jpg, jpeg, webp, gif, bmp, tiff)
        elif filename.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff")):
            ext = filename.split(".")[-1]
            if ext in ["gif", "bmp", "tiff"]:
                img = Image.open(BytesIO(file_bytes))
                png_io = BytesIO()
                img.convert("RGB").save(png_io, format="PNG")
                file_bytes = png_io.getvalue()
                media_type = "image/png"
            else:
                media_type = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                
            encoded_image = base64.b64encode(file_bytes).decode("utf-8")
            pending_attachments.append({
                "filename": file.filename,
                "type": "image",
                "media_type": media_type,
                "data": encoded_image
            })
            return {
                "status": f"Attached image {file.filename}",
                "attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]
            }
            
        # 5. Process Text/Code files
        elif filename.endswith((".txt", ".csv", ".json", ".xml", ".py", ".js", ".html", ".css", ".md", ".yaml", ".yml", ".sh", ".bat", ".sql")):
            try:
                extracted_text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                extracted_text = file_bytes.decode("latin-1", errors="ignore")
                
            pending_attachments.append({
                "filename": file.filename,
                "type": "text",
                "content": extracted_text
            })
            return {
                "status": f"Parsed text file {file.filename}",
                "attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]
            }
        
        else:
            # Fallback: try reading as text
            try:
                extracted_text = file_bytes.decode("utf-8")
                pending_attachments.append({
                    "filename": file.filename,
                    "type": "text",
                    "content": extracted_text
                })
                return {
                    "status": f"Parsed text file {file.filename}",
                    "attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]
                }
            except Exception:
                raise HTTPException(status_code=400, detail="Unsupported file format")

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to process file: {str(e)}")


# --- REMOVE ATTACHMENT ENDPOINT ---
@app.post("/delete_attachment")
async def delete_attachment(request: Request):
    data = await request.json()
    filename = data.get("filename")
    global pending_attachments
    pending_attachments = [a for a in pending_attachments if a["filename"] != filename]
    return {"attachments": [{"filename": a["filename"], "type": a["type"]} for a in pending_attachments]}


# --- CONVERSATIONAL CHAT ROUTER ---
@app.post("/chat")
async def chat(request: Request):
    data = await request.json()
    user_msg = data["message"]
    global pending_attachments
    
    message_content = []
    fallback_message_content = []
    has_pdf = False
    
    for attachment in pending_attachments:
        if attachment["type"] == "text":
            block = {
                "type": "text",
                "text": f"[Uploaded File: {attachment['filename']}]\n{attachment['content']}"
            }
            message_content.append(block)
            fallback_message_content.append(block)
        elif attachment["type"] == "image":
            img_block = {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": attachment["media_type"],
                    "data": attachment["data"]
                }
            }
            text_block = {
                "type": "text",
                "text": f"(Attached image file: {attachment['filename']})"
            }
            message_content.extend([img_block, text_block])
            fallback_message_content.extend([img_block, text_block])
        elif attachment["type"] == "pdf":
            has_pdf = True
            doc_block = {
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": attachment["data"]
                }
            }
            text_desc_block = {
                "type": "text",
                "text": f"(Attached PDF file: {attachment['filename']})"
            }
            message_content.extend([doc_block, text_desc_block])
            
            # Fallback for text
            text_fallback_block = {
                "type": "text",
                "text": f"[Uploaded PDF File: {attachment['filename']}]\n{attachment['content']}"
            }
            fallback_message_content.append(text_fallback_block)

    # Add user query text block
    query_block = {
        "type": "text",
        "text": user_msg
    }
    message_content.append(query_block)
    fallback_message_content.append(query_block)

    response = None
    used_fallback = False
    
    # Route requests according to PDF attachment presence to prevent header errors
    if has_pdf:
        try:
            test_history = history + [{"role": "user", "content": message_content}]
            response = client.messages.create(
                model="claude-fable-5",
                max_tokens=8192,
                system=SYSTEM_PROMPT,
                messages=test_history,
                extra_headers={"anthropic-beta": "pdf-2024-10-22"}
            )
            # If succeeded, append native message content to history
            history.append({"role": "user", "content": message_content})
        except Exception as e:
            # Fallback to sending the extracted text (without beta header)
            try:
                test_history = history + [{"role": "user", "content": fallback_message_content}]
                response = client.messages.create(
                    model="claude-fable-5",
                    max_tokens=8192,
                    system=SYSTEM_PROMPT,
                    messages=test_history
                )
                history.append({"role": "user", "content": fallback_message_content})
                used_fallback = True
            except Exception as inner_e:
                raise HTTPException(status_code=500, detail=f"Failed to communicate with Claude (retry): {str(inner_e)}")
    else:
        # Standard query: send normally without PDF beta headers
        try:
            test_history = history + [{"role": "user", "content": message_content}]
            response = client.messages.create(
                model="claude-fable-5",
                max_tokens=8192,
                system=SYSTEM_PROMPT,
                messages=test_history
            )
            history.append({"role": "user", "content": message_content})
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to communicate with Claude: {str(e)}")

    pending_attachments.clear()

    reply = ""
    for block in response.content:
        if hasattr(block, "type"):
            if block.type == "text":
                reply += block.text
        elif isinstance(block, dict):
            if block.get("type") == "text":
                reply += block.get("text", "")
    history.append({"role": "assistant", "content": reply})

    downloads = []
    preview = None
    
    # 1. XML / SVG Block
    svg_match = re.search(r"```(xml|svg)\s*([\s\S]*?)(?:```|$)", reply, re.IGNORECASE)
    if svg_match:
        content = svg_match.group(2).strip()
        if "<svg" in content:
            try:
                with open("output.svg", "w", encoding="utf-8") as f:
                    f.write(content)
                downloads.append({"label": "Download SVG Vector", "url": "/download/svg", "type": "svg"})
                preview = content
            except Exception:
                pass

    # 2. HTML -> PDF Block
    html_match = re.search(r"```html\s*([\s\S]*?)(?:```|$)", reply, re.IGNORECASE)
    if html_match:
        content = html_match.group(1).strip()
        try:
            with open("temp.html", "w", encoding="utf-8") as f:
                f.write(content)
            success = convert_html_to_pdf(content, "output.pdf")
            if success:
                downloads.append({"label": "Download PDF Document", "url": "/download/pdf", "type": "pdf"})
        except Exception:
            pass

    # 3. Markdown -> DOCX Block
    md_match = re.search(r"```markdown\s*([\s\S]*?)(?:```|$)", reply, re.IGNORECASE)
    if md_match:
        content = md_match.group(1).strip()
        try:
            doc = docx.Document()
            write_markdown_to_docx(doc, content)
            doc.save("output.docx")
            downloads.append({"label": "Download Word Document", "url": "/download/docx", "type": "docx"})
        except Exception:
            pass

    # 4. JSON Slides or text slides -> PPTX Block
    json_match = re.search(r"```json\s*([\s\S]*?)(?:```|$)", reply, re.IGNORECASE)
    slides_data = None
    if json_match:
        try:
            slides_data = json.loads(json_match.group(1).strip())
            # ensure it's structured slides
            if not (isinstance(slides_data, list) and len(slides_data) > 0 and isinstance(slides_data[0], dict)):
                slides_data = None
        except Exception:
            pass

    if not slides_data:
        # Check if text slides exist
        if any(indicator in reply.lower() for indicator in ["slide 1", "slide 2", "--- slide"]):
            try:
                slides_data = extract_slides_from_text(reply)
            except Exception:
                pass

    if slides_data:
        try:
            prs = Presentation()
            bullet_layout = prs.slide_layouts[1]
            for slide_info in slides_data:
                slide = prs.slides.add_slide(bullet_layout)
                slide.shapes.title.text = slide_info.get("title", "Slide")
                if "content" in slide_info and slide_info["content"]:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    content_list = slide_info["content"]
                    if isinstance(content_list, list):
                        for point in content_list:
                            p = tf.add_paragraph()
                            p.text = str(point)
                            p.level = 0
                    else:
                        p = tf.add_paragraph()
                        p.text = str(content_list)
                        p.level = 0
            prs.save("output.pptx")
            downloads.append({"label": "Download PowerPoint Slides", "url": "/download/pptx", "type": "pptx"})
        except Exception:
            pass

    # 5. Tables -> XLSX Spreadsheet
    if "|" in reply:
        lines = reply.split("\n")
        has_table = False
        for idx in range(len(lines) - 1):
            if lines[idx].strip().startswith("|") and re.match(r"^\|[\s\-\|:]+\|$", lines[idx+1].strip()):
                has_table = True
                break
        if has_table:
            try:
                import openpyxl
                wb = openpyxl.Workbook()
                wb.remove(wb.active) # Remove default sheet
                
                table_count = 0
                i = 0
                while i < len(lines):
                    line = lines[i].strip()
                    if line.startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s\-\|:]+\|$", lines[i+1].strip()):
                        table_count += 1
                        ws = wb.create_sheet(title=f"Table {table_count}")
                        header_cells = [cell.strip() for cell in line.split("|")[1:-1]]
                        ws.append(header_cells)
                        i += 2
                        while i < len(lines) and lines[i].strip().startswith("|"):
                            row_cells = [cell.strip() for cell in lines[i].strip().split("|")[1:-1]]
                            ws.append(row_cells)
                            i += 1
                    else:
                        i += 1
                if table_count > 0:
                    wb.save("output.xlsx")
                    downloads.append({"label": "Download Excel Spreadsheet", "url": "/download/xlsx", "type": "xlsx"})
            except Exception:
                pass

    response_data = {"reply": reply}
    if downloads:
        response_data["downloads"] = downloads
        response_data["download_url"] = downloads[0]["url"]
    if preview:
        response_data["preview"] = preview
        
    return response_data


# --- ON-DEMAND EXPORT ENDPOINTS ---
@app.post("/export/pdf")
async def export_pdf(request: Request):
    data = await request.json()
    text = data.get("text", "")
    
    # Check if text contains an HTML code block
    html_match = re.search(r"```html\s*([\s\S]*?)(?:```|$)", text, re.IGNORECASE)
    if html_match:
        html_content = html_match.group(1).strip()
    else:
        html_content = markdown_to_html(text)
        
    export_path = "export.pdf"
    success = convert_html_to_pdf(html_content, export_path)
    if not success or not os.path.exists(export_path):
        raise HTTPException(status_code=500, detail="Failed to compile PDF")
        
    return FileResponse(export_path, media_type="application/pdf", filename="exported_document.pdf")


@app.post("/export/docx")
async def export_docx(request: Request):
    data = await request.json()
    text = data.get("text", "")
    
    # Check if text contains a markdown code block
    md_match = re.search(r"```markdown\s*([\s\S]*?)(?:```|$)", text, re.IGNORECASE)
    if md_match:
        content = md_match.group(1).strip()
    else:
        content = text
        
    doc = docx.Document()
    write_markdown_to_docx(doc, content)
            
    export_path = "export.docx"
    doc.save(export_path)
    return FileResponse(export_path, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="exported_document.docx")


@app.post("/export/pptx")
async def export_pptx(request: Request):
    data = await request.json()
    text = data.get("text", "")
    
    slides_data = extract_slides_from_text(text)
    prs = Presentation()
    bullet_layout = prs.slide_layouts[1]
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_info.get("title", "Slide")
        
        if "content" in slide_info and slide_info["content"]:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for point in slide_info["content"]:
                p = tf.add_paragraph()
                p.text = str(point)
                p.level = 0
                
    export_path = "export.pptx"
    prs.save(export_path)
    return FileResponse(export_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="exported_presentation.pptx")


@app.post("/export/xlsx")
async def export_xlsx(request: Request):
    data = await request.json()
    text = data.get("text", "")
    
    import openpyxl
    wb = openpyxl.Workbook()
    wb.remove(wb.active) # Remove default sheet
    
    lines = text.split("\n")
    table_count = 0
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if line.startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s\-\|:]+\|$", lines[i+1].strip()):
            table_count += 1
            ws = wb.create_sheet(title=f"Table {table_count}")
            header_cells = [cell.strip() for cell in line.split("|")[1:-1]]
            ws.append(header_cells)
            i += 2
            while i < len(lines) and lines[i].strip().startswith("|"):
                row_cells = [cell.strip() for cell in lines[i].strip().split("|")[1:-1]]
                ws.append(row_cells)
                i += 1
        else:
            i += 1
            
    if table_count == 0:
        # If no tables found, let's put the raw text in a single sheet
        ws = wb.create_sheet(title="Content")
        for line in lines:
            ws.append([line])
            
    export_path = "export.xlsx"
    wb.save(export_path)
    return FileResponse(export_path, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", filename="exported_data.xlsx")


# --- FILE SERVING ENDPOINTS ---
@app.get("/download/svg")
async def download_svg():
    return FileResponse("output.svg", media_type="image/svg+xml", filename="vector.svg")

@app.get("/download/pdf")
async def download_pdf():
    return FileResponse("output.pdf", media_type="application/pdf", filename="document.pdf")

@app.get("/download/docx")
async def download_docx():
    return FileResponse("output.docx", media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="document.docx")

@app.get("/download/pptx")
async def download_pptx():
    return FileResponse("output.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="presentation.pptx")

@app.get("/download/xlsx")
async def download_xlsx():
    return FileResponse("output.xlsx", media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", filename="spreadsheet.xlsx")


@app.post("/reset")
async def reset():
    history.clear()
    global pending_attachments
    pending_attachments.clear()
    for f in ["output.svg", "output.pdf", "output.docx", "output.pptx", "output.xlsx", "temp.html", "export.pdf", "export.docx", "export.pptx", "export.xlsx"]:
        if os.path.exists(f):
            os.remove(f)
    return {"status": "cleared"}